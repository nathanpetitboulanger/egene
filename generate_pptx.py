from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Génomique : Fondamentaux et Stratégies"
    subtitle.text = "Egene - Génomique & Bioinformatique\nCours d'Introduction"

    # Slide 2: Génétique vs Génomique
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "1. Génétique vs Génomique"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Génétique : Étude de l'hérédité et des gènes individuels."
    p = tf.add_paragraph()
    p.text = "Génomique : Étude globale de l'ensemble du matériel génétique."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Organisation, structure et interactions."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Régions non codantes et influence environnementale."
    p.level = 1

    # Slide 3: Eucaryotes
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "2. Le Domaine des Eucaryotes"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Présence d'introns et d'exons."
    p = tf.add_paragraph()
    p.text = "Forte proportion de séquences répétées."
    p = tf.add_paragraph()
    p.text = "Compartimentation : Noyau, Mitochondrie, Chloroplaste."

    # Slide 4: Histoire Séquençage
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "3. Histoire du Séquençage"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "1ère Gén : Sanger (Précis, lent, fragments courts)."
    p = tf.add_paragraph()
    p.text = "2ème Gén : NGS - Illumina (Massif, parallèle, bas coût)."
    p = tf.add_paragraph()
    p.text = "3ème Gén : TGS - Long Reads (PacBio, Nanopore)."

    # Slide 5: Stratégies
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "4. Stratégies de Séquençage"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "BAC-to-BAC : Hiérarchique, robuste mais lent."
    p = tf.add_paragraph()
    p.text = "Whole Genome Shotgun (WGS) : Rapide, standard actuel."

    # Slide 6: Qualité
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "5. Qualité de l'Assemblage"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "N50 : Métrique de continuité."
    p = tf.add_paragraph()
    p.text = "BUSCO : Métrique de complétude (gènes universels)."
    p = tf.add_paragraph()
    p.text = "Contig (sans lacune) vs Scaffold (avec gaps)."

    # Slide 7: Annotation
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "6. Annotation des Séquences"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Structurale : Identification gènes/exons (Transcriptomique)."
    p = tf.add_paragraph()
    p.text = "Fonctionnelle : Attribution d'une fonction (Homologie)."

    # Slide 8: Bulk vs Single-cell
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "7. Bulk vs Single-cell"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Bulk : Moyenne d'un mélange de cellules."
    p = tf.add_paragraph()
    p.text = "Single-cell : Hétérogénéité et types cellulaires rares."

    # Slide 9: Complexité
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "8. La Complexité du Vivant"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Humain : 3,2 Gb, ~20 000 gènes."
    p = tf.add_paragraph()
    p.text = "Blé : 17 Gb, Hexaploïde, >85% répétitions."
    p = tf.add_paragraph()
    p.text = "Vanille : 4 Gb, Endoréplication partielle."

    # Slide 10: Pangénomique
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "9. La Pangénomique"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Core Genome : Commun à tous (fonctions vitales)."
    p = tf.add_paragraph()
    p.text = "Accessory Genome : Variable (adaptation)."

    # Save
    prs.save('cours/cour_J1_intro/presentation_genomique.pptx')
    print("Présentation générée : cours/cour_J1_intro/presentation_genomique.pptx")

if __name__ == "__main__":
    create_presentation()
